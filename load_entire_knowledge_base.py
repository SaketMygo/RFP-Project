import os
import pandas as pd
from docx import Document
from pypdf import PdfReader
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter
import chromadb
from chromadb.utils import embedding_functions

# ---------------------------------------------------------------------------
# 1. SETUP LOCAL VECTOR DATABASE (CHROMADB)
# ---------------------------------------------------------------------------
# Setting this up will cleanly initialize a new database in your folder
chroma_client = chromadb.PersistentClient(path="./chroma_db")

hf_embedding_fn = embedding_functions.SentenceTransformerEmbeddingFunction(
    model_name="all-MiniLM-L6-v2"
)

collection = chroma_client.get_or_create_collection(
    name="sap_bids_collection", 
    embedding_function=hf_embedding_fn
)

# ---------------------------------------------------------------------------
# 2. PARSING FUNCTIONS FOR ALL ROBUST FORMATS
# ---------------------------------------------------------------------------

def parse_docx(file_path):
    """Extracts text from Word documents."""
    try:
        doc = Document(file_path)
        content = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(content)
    except Exception as e:
        return ""

def parse_pdf(file_path):
    """Extracts text from PDFs page by page."""
    try:
        reader = PdfReader(file_path)
        content = [page.extract_text() for page in reader.pages if page.extract_text()]
        return "\n\n".join(content)
    except Exception as e:
        return ""

def parse_pptx(file_path):
    """Extracts text content from PowerPoint presentation slides."""
    try:
        prs = Presentation(file_path)
        slide_texts = []
        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text.strip())
            if text_runs:
                slide_texts.append(f"--- Slide {i+1} ---\n" + "\n".join(text_runs))
        return "\n\n".join(slide_texts)
    except Exception as e:
        return ""

def parse_excel_to_chunks(file_path, file_name, relative_folder):
    """Converts structured Excel spreadsheet rows into meaningful text chunks."""
    chunks = []
    try:
        xl = pd.ExcelFile(file_path)
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            df.dropna(how='all', inplace=True)
            if df.empty:
                continue
                
            headers = [str(col).strip() for col in df.columns]
            
            for idx, row in df.iterrows():
                row_values = [str(val).strip() for val in row.values]
                row_items = [f"{headers[i]}: {row_values[i]}" for i in range(len(headers)) if row_values[i] != 'nan']
                row_string = ", ".join(row_items)
                
                chunk_text = f"Category: {relative_folder} | Document: {file_name} | Sheet: {sheet_name} | Row {idx+1}: {row_string}"
                
                metadata = {
                    "source_file": file_name,
                    "folder_path": relative_folder,
                    "document_type": "RFQ/Excel",
                    "sheet_name": sheet_name
                }
                chunks.append({"text": chunk_text, "metadata": metadata})
    except Exception as e:
        pass
    return chunks

# ---------------------------------------------------------------------------
# 3. COMPREHENSIVE FOLDER SCANNING LOOP
# ---------------------------------------------------------------------------

def ingest_master_directory(root_dir):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=120)
    
    global_id_counter = 0
    total_files_processed = 0

    print(f"🚀 Initializing full sweep of Knowledge Base: {root_dir}\n")

    for dirpath, dirnames, filenames in os.walk(root_dir):
        # Captures top-level grouping (e.g., Active Bids, Close Won) and nested structures
        relative_folder = os.path.relpath(dirpath, root_dir)
        if relative_folder == ".":
            relative_folder = "Root"

        for file_name in filenames:
            # Skip hidden temporary/ghost files standard to MS Office updates
            if file_name.startswith("~$"):
                continue
                
            file_path = os.path.join(dirpath, file_name)
            ext = os.path.splitext(file_name)[1].lower()
            
            # Infer precise document category tag
            doc_type = "Template/General"
            lower_path = file_path.lower()
            if "active bids" in lower_path: doc_type = "Active Bid"
            elif "close won" in lower_path: doc_type = "Closed Won RFP"
            elif "lost or no bid" in lower_path: doc_type = "Closed Lost RFP"
            elif "open" in lower_path: doc_type = "Open Opportunity"

            base_metadata = {
                "source_file": file_name,
                "folder_path": relative_folder,
                "document_type": doc_type
            }

            file_texts = []
            file_metadatas = []
            file_ids = []

            # Word Docs
            if ext == '.docx':
                print(f"Parsing Word  📄 [{relative_folder}] -> {file_name}")
                raw_text = parse_docx(file_path)
                if raw_text:
                    chunks = text_splitter.split_text(raw_text)
                    for chunk in chunks:
                        file_texts.append(chunk)
                        file_metadatas.append(base_metadata)
                        file_ids.append(f"doc_{global_id_counter}")
                        global_id_counter += 1

            # PDFs
            elif ext == '.pdf':
                print(f"Parsing PDF   📕 [{relative_folder}] -> {file_name}")
                raw_text = parse_pdf(file_path)
                if raw_text:
                    chunks = text_splitter.split_text(raw_text)
                    for chunk in chunks:
                        file_texts.append(chunk)
                        file_metadatas.append(base_metadata)
                        file_ids.append(f"doc_{global_id_counter}")
                        global_id_counter += 1

            # Presentations
            elif ext in ['.pptx', '.ppt']:
                print(f"Parsing Slides📙 [{relative_folder}] -> {file_name}")
                raw_text = parse_pptx(file_path)
                if raw_text:
                    chunks = text_splitter.split_text(raw_text)
                    for chunk in chunks:
                        file_texts.append(chunk)
                        file_metadatas.append(base_metadata)
                        file_ids.append(f"doc_{global_id_counter}")
                        global_id_counter += 1

            # Spreadsheets
            elif ext in ['.xlsx', '.xls']:
                print(f"Parsing Excel 📊 [{relative_folder}] -> {file_name}")
                excel_chunks = parse_excel_to_chunks(file_path, file_name, relative_folder)
                for chunk in excel_chunks:
                    if chunk["metadata"]["document_type"] == "Template/General":
                        chunk["metadata"]["document_type"] = doc_type
                    file_texts.append(chunk["text"])
                    file_metadatas.append(chunk["metadata"])
                    file_ids.append(f"doc_{global_id_counter}")
                    global_id_counter += 1

            # Write batches to database
            if file_texts:
                collection.upsert(
                    documents=file_texts,
                    metadatas=file_metadatas,
                    ids=file_ids
                )
                total_files_processed += 1

    print(f"\n✨ System Sync Complete! Indexed {total_files_processed} total files across all operations.")
    print(f"📁 Database now contains {global_id_counter} unified search vectors.")

# ---------------------------------------------------------------------------
# RUN MAIN SCRIPT
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    # Pointing directly to your master directory folder
    master_directory = r"C:\Users\Saket Dronamraju\Desktop\RFP project\01-RFPs-RFIs-RFQs"
    
    if os.path.exists(master_directory):
        ingest_directory_to_chroma(master_directory) if 'ingest_directory_to_chroma' in globals() else ingest_master_directory(master_directory)
    else:
        print(f"❌ Target folder missing: {master_directory}")
